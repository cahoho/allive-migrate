"""
Generate PPT for Tencent Cloud Hackathon - 众生迁徙 (Allive Migration)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color theme - Nature & Ecology
PRIMARY_GREEN = RGBColor(0x1B, 0x5E, 0x20)      # Dark forest green
SECONDARY_GREEN = RGBColor(0x4C, 0xAF, 0x50)      # Medium green
LIGHT_GREEN = RGBColor(0xA5, 0xD6, 0xA7)          # Light green
ACCENT_BLUE = RGBColor(0x15, 0x60, 0xC0)          # Sky blue
ACCENT_ORANGE = RGBColor(0xE6, 0x51, 0x00)        # Warm orange
DARK_BG = RGBColor(0x1A, 0x23, 0x2E)              # Dark background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xDC)            # Cream
DARK_TEXT = RGBColor(0x21, 0x27, 0x21)
MEDIUM_TEXT = RGBColor(0x42, 0x48, 0x42)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
GOLD = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ============================================================
# Helper functions
# ============================================================

def add_bg_rect(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle as background"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    """Add a text box with specified properties"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text_box(slide, left, top, width, height, lines):
    """Add text box with multiple styled lines.
    lines: list of (text, font_size, color, bold, alignment)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text, font_size, color, bold = line_data[:4]
        alignment = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bullet_char='● '):
    """Add a bulleted text list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}{item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(8)
    return txBox

def add_section_header(slide, text, subtitle=""):
    """Add a consistent section header bar at top"""
    # Top accent bar
    bar = add_bg_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), SECONDARY_GREEN)
    # Title background
    add_bg_rect(slide, Inches(0), Inches(0.06), prs.slide_width, Inches(1.2), WHITE)
    # Green left accent
    accent = add_bg_rect(slide, Inches(0), Inches(0.06), Inches(0.12), Inches(1.2), SECONDARY_GREEN)
    # Title text
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                 text, font_size=32, color=PRIMARY_GREEN, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.4),
                     subtitle, font_size=16, color=MEDIUM_TEXT)

def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MEDIUM_TEXT,
                 alignment=PP_ALIGN.RIGHT)

def add_icon_circle(slide, left, top, size, color, label):
    """Add a circle with text label"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape

TOTAL_SLIDES = 11

# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Dark nature background
add_bg_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, DARK_BG)

# Overlay gradient effect - top green tint
add_bg_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(3.8),
            RGBColor(0x1B, 0x3A, 0x2A))

# Decorative line
add_bg_rect(slide, Inches(3.5), Inches(3.0), Inches(6.3), Inches(0.04), SECONDARY_GREEN)

# Main title
add_text_box(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
             '众生迁徙', font_size=64, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# English subtitle
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.6),
             'Allive Migration', font_size=30, color=LIGHT_GREEN,
             alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.8),
             '一个关于生态保护 · 生命敬畏 · 物种共存的交互式游戏', font_size=20,
             color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# Bottom info
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
             '腾讯云 黑客松开发挑战赛', font_size=22, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
             '2026年6月', font_size=16, color=RGBColor(0x99, 0x99, 0x99),
             alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 项目概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '项目概述', 'Project Overview')
add_page_number(slide, 2, TOTAL_SLIDES)

content = [
    ('众生迁徙 是什么？', 22, PRIMARY_GREEN, True),
    ('', 10, DARK_TEXT, False),
    ('众生迁徙是一款基于 Web 的生态保护题材交互式策略游戏。玩家扮演"自然守护者"，', 17, DARK_TEXT, False),
    ('通过为不同物种规划迁徙路线，帮助它们在日益严峻的生态环境中完成迁徙使命。', 17, DARK_TEXT, False),
    ('', 10, DARK_TEXT, False),
    ('核心特征', 20, PRIMARY_GREEN, True),
]
add_rich_text_box(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(4.5), content)

features = [
    '🎮 8种真实迁徙物种（候鸟、帝王蝶、鲑鱼、角马、海龟等）',
    '🌍 动态生成的生态地图，每局都不同',
    '🏭 人类活动实时影响迁徙廊道',
    '🌪️ 季节变化与自然灾害（暴风、干旱）',
    '📊 生态健康值系统（栖息地退化与恢复）',
    '🧬 物种多样性 / 灭绝机制',
    '🎵 AI生成的沉浸式自然音效',
]
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(7.2), Inches(3.0),
                features, font_size=15, color=DARK_TEXT, bullet_char='')

# Right side: key numbers
add_bg_rect(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(5.2),
            RGBColor(0xE8, 0xF5, 0xE9))

numbers = [
    ('技术栈', 'Vue 3 + TypeScript + SVG + CSS', SECONDARY_GREEN),
    ('代码规模', '~50 源文件, ~5000+ 行代码', ACCENT_BLUE),
    ('物种数量', '8 种迁徙动物', ACCENT_ORANGE),
    ('开发周期', '基于 CodeBuddy AI 全流程驱动', PRIMARY_GREEN),
]
for i, (label, value, color) in enumerate(numbers):
    y = Inches(1.8 + i * 1.2)
    add_text_box(slide, Inches(9.0), y, Inches(3.4), Inches(0.4),
                 label, font_size=12, color=MEDIUM_TEXT, bold=True)
    add_text_box(slide, Inches(9.0), y + Inches(0.3), Inches(3.4), Inches(0.6),
                 value, font_size=15, color=color, bold=True)


# ============================================================
# Slide 3: 游戏理念 - 小红花游戏 + 生态保护
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '游戏理念', 'Game Philosophy — 小红花游戏 × 生态保护')
add_page_number(slide, 3, TOTAL_SLIDES)

# Quote section
add_bg_rect(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.5),
            RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(11), Inches(1.3),
             '「小红花游戏」理念：每一次成功的迁徙，都是一朵"小红花"——\n'
             '它不仅是分数的累积，更是对自然生命的一次守护行动。',
             font_size=18, color=PRIMARY_GREEN, bold=False,
             alignment=PP_ALIGN.LEFT)

# Core philosophy
philosophy = [
    ('🌱 激发敬畏之心', '通过游戏机制，让玩家直观感受物种迁徙的艰难与脆弱，\n激发对生态保护和自然生命的敬畏与呵护之心。'),
    ('🔄 动态生态网络', '每一个节点都是生态链条的一环，人类活动的每一次干预都会在全球生态网络中产生连锁反应。'),
    ('⚖️ 平衡与共存', '游戏不是"消灭人类"，而是展示人类与自然如何找到共存之道——缩小的缩圈点、恢复的栖息地，都是希望的象征。'),
    ('📖 寓教于乐', '每种物种都附带生态学启示文案，让玩家在游戏中学习真实世界的生态知识。'),
]
for i, (title, desc) in enumerate(philosophy):
    y = Inches(3.3 + i * 1.05)
    add_text_box(slide, Inches(0.8), y, Inches(6), Inches(0.35),
                 title, font_size=16, color=PRIMARY_GREEN, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(6.5), Inches(0.65),
                 desc, font_size=13, color=MEDIUM_TEXT)


# ============================================================
# Slide 4: 项目架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '项目架构', 'Project Architecture')
add_page_number(slide, 4, TOTAL_SLIDES)

# Tech Stack
tech_items = [
    ('Vue 3', 'Composition API\n+ script setup', SECONDARY_GREEN),
    ('Vite 5', '构建工具\nHMR 热更新', ACCENT_BLUE),
    ('TypeScript', '严格模式\n类型安全', RGBColor(0x31, 0x78, 0xC6)),
    ('SVG + CSS', '矢量渲染\n无Canvas依赖', ACCENT_ORANGE),
    ('自研引擎', '地图/任务/路线\nreactive store', RGBColor(0x7B, 0x1F, 0xA2)),
]
for i, (name, desc, color) in enumerate(tech_items):
    x = Inches(0.8 + i * 2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, Inches(1.6), Inches(2.2), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Name
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    # Description below shape
    add_text_box(slide, x, Inches(3.4), Inches(2.2), Inches(0.9),
                 desc, font_size=12, color=MEDIUM_TEXT,
                 alignment=PP_ALIGN.CENTER)

# Architecture diagram
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(4), Inches(0.4),
             '🏗️ 核心代码架构', font_size=18, color=PRIMARY_GREEN, bold=True)

arch_items = [
    'src/data/        — 游戏配置、物种模板、事件定义、教程步骤',
    'src/systems/     — 地图生成、任务生成、路线求解、生态健康、音频管理',
    'src/store/       — 全局游戏状态管理 (Vue reactive, ~1900行)',
    'src/components/  — 15个Vue组件 (GameMap, TopBar, SpeciesPanel等)',
    'src/composables/ — 拖线路由、游戏循环、路线验证',
    'src/utils/       — 几何计算、SVG路径工具',
]
add_bullet_list(slide, Inches(0.8), Inches(4.85), Inches(7), Inches(2.3),
                arch_items, font_size=12, color=DARK_TEXT, bullet_char='')


# ============================================================
# Slide 5: 核心系统
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '核心系统设计', 'Core Game Systems')
add_page_number(slide, 5, TOTAL_SLIDES)

systems = [
    ('动态地图\n生成', '开局6节点 → 成功迁徙解锁新节点\n最大24个节点，单一连通分量\n330px自动连边，每节点≤5连接', PRIMARY_GREEN),
    ('动态任务\n系统', '起点/终点随机抽取 + 必经点约束\nBFS/DFS多解保障\n资源消耗 ≤ 14单位', ACCENT_BLUE),
    ('连续难度\n曲线', '阶段1→6渐进解锁\nmaxConcurrent: 1→5\n任务密度随阶段递增', RGBColor(0x7B, 0x1F, 0xA2)),
    ('路线求解\n引擎', '多层校验：节点可用性→边连通→\ntag/必经点→灾种敏感→资源预算\n拒绝无解任务，保证公平', ACCENT_ORANGE),
    ('世界事件\n系统', '5种事件：暴风/野火/干旱/滑坡/繁盛\n三阶段：warning→active→recovery\n公平性预检保证任务可解', RGBColor(0x00, 0x89, 0x7B)),
    ('人类活动\n系统', '唯一永久阻挡圆，鼠标拖拽引导\n吞噬迁徙物种→人类圆变大\n缩圈点机制→玩家可缩小人类影响', RGBColor(0xD4, 0x3F, 0x1A)),
]
for i, (title, desc, color) in enumerate(systems):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.6 + row * 2.9)

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(3.8), Inches(2.6))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)

    # Color accent bar
    add_bg_rect(slide, x, y, Inches(3.8), Inches(0.06), color)

    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.7),
                 title, font_size=16, color=color, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(0.95), Inches(3.4), Inches(1.5),
                 desc, font_size=12, color=MEDIUM_TEXT)


# ============================================================
# Slide 6: AI能力 - CodeBuddy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 'AI 能力应用', 'CodeBuddy — AI 驱动的全流程开发')
add_page_number(slide, 6, TOTAL_SLIDES)

# Main intro
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.9),
             '本项目以 CodeBuddy 为 AI 主力开发工具，覆盖了从架构设计、代码实现、调试优化到\n'
             '文档生成、项目管理的完整开发流程。AI 不仅是一个"辅助工具"，更是核心生产力引擎。',
             font_size=16, color=DARK_TEXT)

# AI Usage categories
ai_categories = [
    ('🎯 架构设计\n与规划', [
        '项目结构设计（data / systems / store / components）',
        '游戏机制数学建模（难度曲线、生态健康公式）',
        'API 接口与数据流设计',
    ]),
    ('💻 代码生成\n与实现', [
        '~50个源文件，5000+行 TypeScript/Vue 代码',
        '复杂算法：BFS/DFS路线求解、地图生成、人类追击预测',
        '组件开发：15个Vue组件，响应式状态管理',
    ]),
    ('🐛 调试\n与优化', [
        '类型检查 + 编译错误修复 (vue-tsc)',
        '运行时逻辑纠错（路线验证、物种多样性判定）',
        '性能优化（避免主线程阻塞、RAF 循环设计）',
    ]),
    ('🎨 多媒体\n内容', [
        'AI 生成 11 种游戏背景音乐（WAV格式）',
        '18 个物种/事件音效',
        'SVG 矢量图标与视觉素材',
    ]),
    ('📝 文档\n与管理', [
        'README / 设计文档 / 开发日志自动生成',
        'CodeBuddy Memory 项目记忆持续积累',
        '版本迭代追踪与策划需求管理',
    ]),
]
for i, (title, items) in enumerate(ai_categories):
    col = i % 3 if i < 3 else i % 3
    row = 0 if i < 3 else 1
    x = Inches(0.6 + col * 4.1)
    y = Inches(2.6 + row * 2.5)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(3.8), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)

    add_bg_rect(slide, x, y, Inches(3.8), Inches(0.06), SECONDARY_GREEN)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.6),
                 title, font_size=15, color=PRIMARY_GREEN, bold=True)

    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.25), y + Inches(0.85 + j * 0.32),
                     Inches(3.3), Inches(0.3),
                     f'• {item}', font_size=11, color=MEDIUM_TEXT)


# ============================================================
# Slide 7: 物种与生态系统
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '物种与生态系统', 'Species & Ecology System')
add_page_number(slide, 7, TOTAL_SLIDES)

# Species grid
species_list = [
    ('🐦 雁鸭候鸟', 'Migratory Waterbird', '1分', '0.55s', '#7AD7FF'),
    ('🦋 帝王蝶', 'Monarch Butterfly', '1.2分', '0.85s', '#FFB45C'),
    ('🦆 斑头雁', 'Bar-headed Goose', '0.9分', '0.42s', '#CFEAFF'),
    ('🐟 鲑鱼', 'Salmon', '1.5分', '1.25s', '#5CA8FF'),
    ('🐃 角马兽群', 'Wildebeest Herd', '3.4分', '7.2s', '#D9C27A'),
    ('🐍 美洲鳗', 'American Eel', '1.7分', '1.65s', '#76D0C4'),
    ('🐢 绿海龟', 'Green Sea Turtle', '2.7分', '4.8s', '#7EDC93'),
    ('🐸 林蛙', 'Wood Frog', '3.8分', '9.2s', '#8CCB5E'),
]
for i, (name, eng, score, speed, color_hex) in enumerate(species_list):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.6 + row * 2.7)
    color = RGBColor(*[int(color_hex[j:j+2], 16) for j in range(1, 7, 2)])

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(2.8), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_bg_rect(slide, x, y, Inches(2.8), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.5),
                 name, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.5), Inches(0.35),
                 eng, font_size=11, color=MEDIUM_TEXT)

    # Info
    add_text_box(slide, x + Inches(0.15), y + Inches(1.1), Inches(2.5), Inches(0.3),
                 f'迁徙得分：{score}', font_size=12, color=DARK_TEXT)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.4), Inches(2.5), Inches(0.3),
                 f'每段速度：{speed}', font_size=12, color=MEDIUM_TEXT)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.7), Inches(2.5), Inches(0.5),
                 f'阶段 {i+1 if i < 2 else (i if i < 3 else i+1 if i < 4 else i+1 if i < 6 else 6)} 解锁',
                 font_size=11, color=RGBColor(0x99, 0x99, 0x99))

# Design principle note
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             '💡 设计原则：越慢/越依赖连续通道的物种 → 保护价值越大 → 迁徙得分越高',
             font_size=11, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 8: 人类活动与生态影响
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '人类活动与生态影响', 'Human Activity & Ecological Impact')
add_page_number(slide, 8, TOTAL_SLIDES)

# Left: Human system
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             '🏭 人类活动系统', font_size=20, color=PRIMARY_GREEN, bold=True)

human_features = [
    '唯一永久阻挡圆：全图单个人类活动圈，模拟人类开发对生态廊道的切割',
    '鼠标拖拽引导：玩家可控制人类活动圈移动方向',
    '吞噬物种：圈覆盖迁徙中动物 → 人类圈变大，物种失败计数+1',
    '缩圈点机制：绿色缩圈点随机刷新，玩家引导人类吞下 → 体型缩小',
    '追击模式：成功迁徙≥12次后，AI主动预测并追击迁徙中动物',
    '屏保模式：无操作时自动游荡，模拟自然人类活动扩散',
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(3.0),
                human_features, font_size=13, color=DARK_TEXT)

# Right: Eco system
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5),
             '🌍 生态健康系统', font_size=20, color=PRIMARY_GREEN, bold=True)

eco_features = [
    '四级生态状态：Healthy → Pressured → Damaged → Degraded',
    '人类覆盖节点 → 每秒扣健康值2.0',
    '动物迁徙经过 → 每个节点扣1.25（承载压力）',
    '无干扰自动恢复：0.45/秒（延迟2.5秒开始）',
    'Health≤0 → 节点枯竭不可用 → 恢复至35后重新通行',
    '生态教育提示：节点状态变化时推送保护理念',
]
add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(3.0),
                eco_features, font_size=13, color=DARK_TEXT)

# Bottom: Species diversity
add_bg_rect(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5),
            RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.3),
             '🧬 物种多样性机制\n'
             '• 同一物种失败3次 → 灭绝（不再生成新任务）\n'
             '• 物种多样性 = 存活物种 / 已解锁物种 × 100%\n'
             '• 多样性 ≤ 33.3% → 生态崩溃，游戏结束（Game Over）',
             font_size=14, color=RGBColor(0xBF, 0x36, 0x0C))


# ============================================================
# Slide 9: 团队
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '团队成员', 'Team')
add_page_number(slide, 9, TOTAL_SLIDES)

# Center the team card
add_bg_rect(slide, Inches(3.5), Inches(2.2), Inches(6.3), Inches(4.2),
            RGBColor(0xE8, 0xF5, 0xE9))

# Name
add_text_box(slide, Inches(3.5), Inches(2.5), Inches(6.3), Inches(0.8),
             'cahoho', font_size=40, color=PRIMARY_GREEN, bold=True,
             alignment=PP_ALIGN.CENTER)

# Role
add_text_box(slide, Inches(3.5), Inches(3.3), Inches(6.3), Inches(0.5),
             '全栈开发者 · 游戏设计师 · AI 应用探索者',
             font_size=18, color=SECONDARY_GREEN,
             alignment=PP_ALIGN.CENTER)

# Separator
add_bg_rect(slide, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.03), SECONDARY_GREEN)

# Skills/tags
tags = ['Vue 3', 'TypeScript', '游戏设计', 'AI 应用', '生态学']
for i, tag in enumerate(tags):
    x = Inches(4.0 + i * 1.2)
    tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, Inches(4.2), Inches(1.0), Inches(0.4))
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = SECONDARY_GREEN
    tag_shape.line.fill.background()
    tf = tag_shape.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

# About
add_text_box(slide, Inches(4.0), Inches(5.0), Inches(5.3), Inches(1.0),
             '独立完成项目从 0 到 1 的全流程：\n'
             '游戏策划 → 架构设计 → 代码实现 → AI 内容生成 → 测试发布',
             font_size=14, color=MEDIUM_TEXT,
             alignment=PP_ALIGN.CENTER)

# AI partner note
add_text_box(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.5),
             '🤖 AI 搭档：CodeBuddy（深度参与全流程开发）',
             font_size=15, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 10: 黑客松亮点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '黑客松项目亮点', 'Hackathon Highlights')
add_page_number(slide, 10, TOTAL_SLIDES)

highlights = [
    ('1', 'AI 全流程驱动', '从架构设计到代码实现、从音效生成到文档管理，\nCodeBuddy AI 深度参与每一个环节，实现了\"AI + 人\"的高效协作范式。'),
    ('2', '纯粹 Web 技术', '零外部素材依赖，纯 Vue 3 + SVG + CSS 构建完整游戏。\n无 Canvas / Phaser / Unity，展现 Web 原生技术极限。'),
    ('3', '深度游戏机制', '8种物种×6个阶段×动态地图×生态健康×人类活动×\n物种多样性→ 形成了一个复杂而自洽的生态模拟系统。'),
    ('4', '社会价值导向', '不仅是游戏，更是生态保护教育工具。每一次交互都在\n传递\"敬畏自然、守护生命\"的理念。'),
    ('5', '一人成军', '单人（cahoho）+ CodeBuddy AI 搭档，在有限时间内\n完成从 0 到完整的游戏产品。'),
    ('6', '可持续迭代', '项目结构清晰、模块化设计，支持持续扩展新物种、\n新事件、新机制。已产出完善的开发文档体系。'),
]
for i, (num, title, desc) in enumerate(highlights):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.6 + row * 1.85)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY_GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.7), y, Inches(5.2), Inches(0.4),
                 title, font_size=18, color=PRIMARY_GREEN, bold=True)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(5.2), Inches(0.8),
                 desc, font_size=12, color=MEDIUM_TEXT)


# ============================================================
# Slide 11: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Dark background
add_bg_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, DARK_BG)
add_bg_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(3.8),
            RGBColor(0x1B, 0x3A, 0x2A))

# Decorative
add_bg_rect(slide, Inches(3.5), Inches(2.8), Inches(6.3), Inches(0.04), SECONDARY_GREEN)

add_text_box(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
             '谢谢！', font_size=56, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10.3), Inches(0.6),
             'Thank You', font_size=28, color=LIGHT_GREEN,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
             '众生迁徙 · Allive Migration\n'
             '让每一次迁徙，都成为对自然生命的守护',
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.CENTER)

add_bg_rect(slide, Inches(5.0), Inches(4.5), Inches(3.3), Inches(0.03), SECONDARY_GREEN)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
             '开发者：cahoho  |  AI 搭档：CodeBuddy',
             font_size=16, color=GOLD,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.5),
             '腾讯云 黑客松开发挑战赛 · 2026',
             font_size=14, color=RGBColor(0x99, 0x99, 0x99),
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), 'output', '众生迁徙_黑客松PPT.pptx')
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
